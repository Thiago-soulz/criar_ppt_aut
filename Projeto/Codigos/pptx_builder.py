import copy
import os
import tempfile

from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from config import (
    TEMPLATE_PATH,
    OUTPUT_PATH,
    PAGINA_GRID_X_INICIO, PAGINA_GRID_Y_INICIO, PAGINA_GRID_MAX_COLUNAS,
    PAGINA_ESPACAMENTO_X, PAGINA_ESPACAMENTO_Y,
    PAGINA_IMAGEM_LARGURA, PAGINA_IMAGEM_ALTURA,
    RECORTE_GRID_X_INICIO, RECORTE_GRID_Y_INICIO, RECORTE_GRID_MAX_COLUNAS,
    RECORTE_ESPACAMENTO_X, RECORTE_ESPACAMENTO_Y,
    RECORTE_IMAGEM_LARGURA, RECORTE_IMAGEM_ALTURA,
    TITULO_COR_FUNDO, TITULO_COR_TEXTO, TITULO_FONTE_TAMANHO,
    TITULO_X, TITULO_Y, TITULO_LARGURA, TITULO_ALTURA,
)


def _gerar_caminho_output(nome_excel):
    """Gera o caminho do .pptx com base no nome do Excel.
    Se ja existir, incrementa: Arquivo.pptx -> Arquivo.1.pptx -> Arquivo.2.pptx ..."""
    base_dir = os.path.dirname(OUTPUT_PATH)
    nome_base = os.path.splitext(os.path.basename(nome_excel))[0]
    caminho = os.path.join(base_dir, f"{nome_base}.pptx")
    if not os.path.exists(caminho):
        return caminho
    contador = 1
    while True:
        caminho = os.path.join(base_dir, f"{nome_base}.{contador}.pptx")
        if not os.path.exists(caminho):
            return caminho
        contador += 1


def montar_apresentacao(pallets_pdfs, modo, nome_excel=None, cfg_override=None, callback_status=None):

    # Usa cfg_override se a pessoa editou as configurações, senão usa config.py
    if cfg_override:
        cfg = cfg_override
    elif modo == "pagina":
        cfg = {
            "x_inicio":    PAGINA_GRID_X_INICIO,
            "y_inicio":    PAGINA_GRID_Y_INICIO,
            "colunas":     PAGINA_GRID_MAX_COLUNAS,
            "esp_x":       PAGINA_ESPACAMENTO_X,
            "esp_y":       PAGINA_ESPACAMENTO_Y,
            "largura_max": PAGINA_IMAGEM_LARGURA,
            "altura_max":  PAGINA_IMAGEM_ALTURA,
        }
    else:
        cfg = {
            "x_inicio":    RECORTE_GRID_X_INICIO,
            "y_inicio":    RECORTE_GRID_Y_INICIO,
            "colunas":     RECORTE_GRID_MAX_COLUNAS,
            "esp_x":       RECORTE_ESPACAMENTO_X,
            "esp_y":       RECORTE_ESPACAMENTO_Y,
            "largura_max": RECORTE_IMAGEM_LARGURA,
            "altura_max":  RECORTE_IMAGEM_ALTURA,
        }

    prs = Presentation(TEMPLATE_PATH)

    SLIDE_BASE = 1

    for numero_pallet, arquivos_pdf in pallets_pdfs.items():

        slide = _duplicar_slide(prs, SLIDE_BASE)
        _mover_slide(prs, len(prs.slides) - 1, len(prs.slides) - 2)
        _adicionar_titulo(slide, f"Pallet {numero_pallet:02d}")
        _adicionar_imagens(slide, arquivos_pdf, cfg, callback_status)

        if callback_status:
            callback_status(f"Pallet {numero_pallet:02d} finalizado")

    _remover_slide(prs, SLIDE_BASE)

    output_path = _gerar_caminho_output(nome_excel) if nome_excel else OUTPUT_PATH
    prs.save(output_path)

    return output_path


def _duplicar_slide(prs, index):

    template_slide = prs.slides[index]
    layout = template_slide.slide_layout

    novo_slide = prs.slides.add_slide(layout)

    # Remove placeholders que vêm do layout
    for shape in list(novo_slide.placeholders):
        shape.element.getparent().remove(shape.element)

    # Copia só os shapes de conteúdo do slide base
    spTree_novo = novo_slide.shapes._spTree
    spTree_orig = template_slide.shapes._spTree

    for elem in spTree_orig:
        if elem.tag == qn('p:sp'):
            spTree_novo.append(copy.deepcopy(elem))

    return novo_slide


def _adicionar_imagens(slide, arquivos_pdf, cfg, callback_status):

    x = cfg["x_inicio"]
    y = cfg["y_inicio"]
    coluna = 0
    imagens_temp = []

    for pdf_path in arquivos_pdf:

        nome_pdf = os.path.basename(pdf_path)

        if callback_status:
            callback_status(f"Processando {nome_pdf}")

        nome_sem_ext = os.path.splitext(nome_pdf)[0]
        imagem_temp = os.path.join(tempfile.gettempdir(), f"{nome_sem_ext}.png")
        imagens_temp.append(imagem_temp)

        try:

            if not os.path.exists(imagem_temp) or os.path.getsize(imagem_temp) == 0:
                continue

            img = Image.open(imagem_temp)
            largura_img, altura_img = img.size
            proporcao = largura_img / altura_img

            altura_slide = cfg["altura_max"]
            largura_slide = altura_slide * proporcao

            if largura_slide > cfg["largura_max"]:
                largura_slide = cfg["largura_max"]
                altura_slide = largura_slide / proporcao

            slide.shapes.add_picture(
                imagem_temp,
                Inches(x),
                Inches(y),
                width=Inches(largura_slide),
                height=Inches(altura_slide)
            )

            x += largura_slide + cfg["esp_x"]
            coluna += 1

            if coluna >= cfg["colunas"]:
                coluna = 0
                x = cfg["x_inicio"]
                y += altura_slide + cfg["esp_y"]

        except Exception as erro:
            print(f"Erro ao inserir imagem {nome_pdf}: {erro}")

    for imagem in imagens_temp:
        try:
            os.remove(imagem)
        except Exception:
            pass


def _adicionar_titulo(slide, texto):

    caixa = slide.shapes.add_textbox(
        Inches(TITULO_X), Inches(TITULO_Y),
        Inches(TITULO_LARGURA), Inches(TITULO_ALTURA)
    )

    caixa.fill.solid()
    caixa.fill.fore_color.rgb = RGBColor(*TITULO_COR_FUNDO)
    caixa.line.color.rgb = RGBColor(*TITULO_COR_FUNDO)

    frame = caixa.text_frame
    frame.clear()

    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(TITULO_FONTE_TAMANHO)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITULO_COR_TEXTO)


def _mover_slide(prs, old_index, new_index):

    slides = prs.slides._sldIdLst
    slide = slides[old_index]
    slides.remove(slide)
    slides.insert(new_index, slide)


def _remover_slide(prs, index):

    slides = prs.slides._sldIdLst
    slide = list(slides)[index]
    slides.remove(slide)